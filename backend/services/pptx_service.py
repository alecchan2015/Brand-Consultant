"""
Professional PPTX Renderer (v2) — McKinsey-grade slide generation.

Entry point:
    await generate_pptx(task_id, agent_type, brand_name, content, file_path, db)

Flow:
    1. pptx_agent.generate_slide_structure() → structured JSON via LLM
    2. Fallback: _parse_markdown_structure()  → JSON from raw markdown
    3. _render_structure()                    → python-pptx PPTX file
"""
from __future__ import annotations

import re
from datetime import datetime
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Slide geometry ─────────────────────────────────────────────────────────────
SW   = Inches(13.33)   # 16:9 width
SH   = Inches(7.50)    # 16:9 height
MX   = Inches(0.55)    # horizontal margin
HDR_H = Inches(0.95)   # header zone height
FTR_Y = Inches(7.12)   # footer top edge
FTR_H = Inches(0.38)   # footer height
CONT_Y = HDR_H
CONT_H = FTR_Y - CONT_Y
CONT_W = SW - 2 * MX

FONT = "微软雅黑"

# ── Per-agent colour themes ────────────────────────────────────────────────────
_THEMES: Dict[str, Dict[str, RGBColor]] = {
    "strategy": dict(
        primary   = RGBColor(0x1B, 0x2A, 0x4A),
        accent    = RGBColor(0xE0, 0x35, 0x3A),
        secondary = RGBColor(0x2E, 0x75, 0xB6),
        light     = RGBColor(0xF2, 0xF5, 0xF9),
        text      = RGBColor(0x1A, 0x1A, 0x2E),
        white     = RGBColor(0xFF, 0xFF, 0xFF),
        gray      = RGBColor(0x8E, 0x9C, 0xB0),
        muted     = RGBColor(0xBF, 0xCA, 0xDD),
        s_clr     = RGBColor(0x1A, 0xBC, 0x9C),
        w_clr     = RGBColor(0xE7, 0x4C, 0x3C),
        o_clr     = RGBColor(0x29, 0x80, 0xB9),
        t_clr     = RGBColor(0xF3, 0x9C, 0x12),
    ),
    "brand": dict(
        primary   = RGBColor(0x2D, 0x1B, 0x69),
        accent    = RGBColor(0xF5, 0xA6, 0x23),
        secondary = RGBColor(0x6C, 0x5C, 0xE7),
        light     = RGBColor(0xF8, 0xF5, 0xFF),
        text      = RGBColor(0x1A, 0x1A, 0x2E),
        white     = RGBColor(0xFF, 0xFF, 0xFF),
        gray      = RGBColor(0x8E, 0x9C, 0xB0),
        muted     = RGBColor(0xC8, 0xBF, 0xE8),
        s_clr     = RGBColor(0x1A, 0xBC, 0x9C),
        w_clr     = RGBColor(0xE7, 0x4C, 0x3C),
        o_clr     = RGBColor(0x29, 0x80, 0xB9),
        t_clr     = RGBColor(0xF3, 0x9C, 0x12),
    ),
    "operations": dict(
        primary   = RGBColor(0x0F, 0x3D, 0x2D),
        accent    = RGBColor(0xFF, 0x6B, 0x35),
        secondary = RGBColor(0x00, 0xA8, 0x7E),
        light     = RGBColor(0xF0, 0xF9, 0xF5),
        text      = RGBColor(0x1A, 0x1A, 0x2E),
        white     = RGBColor(0xFF, 0xFF, 0xFF),
        gray      = RGBColor(0x8E, 0x9C, 0xB0),
        muted     = RGBColor(0xBF, 0xDE, 0xD4),
        s_clr     = RGBColor(0x1A, 0xBC, 0x9C),
        w_clr     = RGBColor(0xE7, 0x4C, 0x3C),
        o_clr     = RGBColor(0x29, 0x80, 0xB9),
        t_clr     = RGBColor(0xF3, 0x9C, 0x12),
    ),
}

def _th(agent_type: str) -> Dict[str, RGBColor]:
    return _THEMES.get(agent_type, _THEMES["strategy"])


# ── Low-level drawing helpers ──────────────────────────────────────────────────

def _rect(slide, l, t, w, h, fill: RGBColor,
          border: Optional[RGBColor] = None, bw: float = 0,
          rounded: bool = False):
    s = slide.shapes.add_shape(5 if rounded else 1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw or 1.0)
    else:
        s.line.fill.background()
    return s


def _oval(slide, l, t, w, h, fill: RGBColor,
          border: Optional[RGBColor] = None):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(2)
    else:
        s.line.fill.background()
    return s


def _line(slide, x1, y1, x2, y2, color: RGBColor, width_pt: float = 1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width_pt)
    return c


def _diamond(slide, l, t, w, h, fill: RGBColor,
             border: Optional[RGBColor] = None, bw: float = 0):
    """Diamond shape (decision node in flowcharts)."""
    s = slide.shapes.add_shape(4, l, t, w, h)   # shape 4 = MSO diamond
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw or 1.5)
    else:
        s.line.fill.background()
    return s


def _hexagon(slide, l, t, w, h, fill: RGBColor):
    """Hexagon shape — used as visual step number holders."""
    s = slide.shapes.add_shape(10, l, t, w, h)  # shape 10 = hexagon
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _lighten(c: RGBColor, f: float) -> RGBColor:
    """Blend colour c toward white by factor f (0=unchanged, 1=white)."""
    return RGBColor(
        int(c[0] + (255 - c[0]) * f),
        int(c[1] + (255 - c[1]) * f),
        int(c[2] + (255 - c[2]) * f),
    )


def _decor_bg(slide, t: dict):
    """
    Subtle decorative background circles — drawn first so they stay behind
    all content.  Creates a professional 'depth' without distraction.
    """
    # Large faint circle, top-right corner
    r1 = Inches(4.0)
    _oval(slide, SW - r1 * 0.55, -r1 * 0.45, r1, r1,
          _lighten(t["primary"], 0.93))
    # Medium accent circle, bottom-left
    r2 = Inches(2.2)
    _oval(slide, -r2 * 0.35, SH - r2 * 0.7, r2, r2,
          _lighten(t["accent"], 0.91))
    # Small secondary dot, mid-right
    r3 = Inches(0.8)
    _oval(slide, SW - Inches(1.2), SH / 2 - r3 / 2, r3, r3,
          _lighten(t["secondary"], 0.88))


def _txb(slide, text: str, l, t, w, h,
         size: float = 11, bold: bool = False, italic: bool = False,
         color: Optional[RGBColor] = None,
         align=PP_ALIGN.LEFT, wrap: bool = True,
         v_anchor=MSO_ANCHOR.TOP):
    """Add a text box containing one paragraph / run."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def _add_run(paragraph, text: str, size: float, bold=False,
             color: Optional[RGBColor] = None):
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return run


def _bullets_txb(slide, items: List[str], l, t, w, h,
                 size: float = 10.5, color: Optional[RGBColor] = None,
                 bullet_color: Optional[RGBColor] = None,
                 line_spacing_pt: float = 6):
    """Render a bulleted list inside a text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(line_spacing_pt if i > 0 else 0)
        # Bullet marker
        r_mark = p.add_run()
        r_mark.text = "▪ "
        r_mark.font.name = FONT
        r_mark.font.size = Pt(size - 0.5)
        if bullet_color:
            r_mark.font.color.rgb = bullet_color
        elif color:
            r_mark.font.color.rgb = color
        # Text
        r_text = p.add_run()
        r_text.text = item
        r_text.font.name = FONT
        r_text.font.size = Pt(size)
        if color:
            r_text.font.color.rgb = color
    return tb


# ── Slide-level chrome (header + footer) ──────────────────────────────────────

def _header(slide, title: str, subtitle: str, t: dict):
    """Thin accent stripe + light bg + title + optional subtitle."""
    # Accent top stripe
    _rect(slide, 0, 0, SW, Inches(0.065), t["accent"])
    # Light header background
    _rect(slide, 0, Inches(0.065), SW, Inches(0.885), t["light"])
    # Vertical left accent bar
    _rect(slide, 0, Inches(0.065), Inches(0.055), Inches(0.885), t["secondary"])

    title_h = Inches(0.52) if subtitle else Inches(0.75)
    _txb(slide, title,
         MX, Inches(0.1), CONT_W, title_h,
         size=22, bold=True, color=t["primary"],
         v_anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _txb(slide, subtitle,
             MX, Inches(0.6), CONT_W, Inches(0.32),
             size=9.5, color=t["gray"])


def _footer(slide, brand_name: str, slide_num: int, t: dict):
    _rect(slide, 0, FTR_Y, SW, Inches(0.02), t["gray"])
    _txb(slide, brand_name,
         MX, FTR_Y + Inches(0.05), Inches(5), Inches(0.28),
         size=8, color=t["gray"])
    _txb(slide, datetime.now().strftime("%Y.%m"),
         SW / 2 - Inches(1), FTR_Y + Inches(0.05), Inches(2), Inches(0.28),
         size=8, color=t["gray"], align=PP_ALIGN.CENTER)
    _txb(slide, str(slide_num),
         SW - MX - Inches(0.6), FTR_Y + Inches(0.05), Inches(0.55), Inches(0.28),
         size=8, bold=True, color=t["accent"], align=PP_ALIGN.RIGHT)


def _blank(prs) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Individual slide renderers ─────────────────────────────────────────────────

def _render_cover(prs, data: dict, t: dict, brand_name: str, _num: int):
    slide = _blank(prs)
    # Dark background
    _rect(slide, 0, 0, SW, SH, t["primary"])
    # Top accent bar
    _rect(slide, 0, 0, SW, Inches(0.14), t["accent"])
    # Left decorative bar
    _rect(slide, 0, 0, Inches(0.35), SH, t["secondary"])
    # Bottom accent bar
    _rect(slide, 0, SH - Inches(0.1), SW, Inches(0.1), t["accent"])

    # Brand circle (top-left decoration)
    _oval(slide,
          Inches(0.65), Inches(0.4), Inches(0.7), Inches(0.7),
          t["accent"])
    initial = (brand_name or "B")[0].upper()
    _txb(slide, initial,
         Inches(0.65), Inches(0.38), Inches(0.7), Inches(0.72),
         size=22, bold=True, color=t["white"], align=PP_ALIGN.CENTER,
         v_anchor=MSO_ANCHOR.MIDDLE)

    # Main title
    title = data.get("title", brand_name)
    _txb(slide, title,
         Inches(1.5), Inches(1.9), Inches(10.5), Inches(1.6),
         size=40, bold=True, color=t["white"],
         align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # Divider
    _rect(slide, Inches(3.5), Inches(3.7), Inches(6.3), Inches(0.05), t["accent"])

    # Subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        _txb(slide, subtitle,
             Inches(1.5), Inches(3.85), Inches(10.5), Inches(0.85),
             size=16, color=t["muted"],
             align=PP_ALIGN.CENTER)

    # Bottom meta
    _txb(slide, brand_name,
         Inches(0.65), SH - Inches(0.7), Inches(4), Inches(0.35),
         size=9.5, color=t["muted"])
    _txb(slide, datetime.now().strftime("%Y年%m月"),
         SW - Inches(2.4), SH - Inches(0.7), Inches(1.85), Inches(0.35),
         size=9.5, color=t["muted"], align=PP_ALIGN.RIGHT)


def _render_agenda(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "报告结构"), "", t)
    _footer(slide, brand_name, num, t)

    items = data.get("items", [])
    n = len(items)
    if not n:
        return

    # Two-column layout for agenda
    col = 2
    rows = (n + col - 1) // col
    card_w = Inches(5.6)
    card_h = Inches(0.75)
    gap_x  = Inches(0.55)
    gap_y  = Inches(0.22)
    start_x = [MX, MX + card_w + gap_x]
    start_y = CONT_Y + Inches(0.3)

    for idx, item in enumerate(items):
        c = idx % col
        r = idx // col
        x = start_x[c]
        y = start_y + r * (card_h + gap_y)

        # Numbered card
        num_w = Inches(0.55)
        _rect(slide, x, y, num_w, card_h, t["primary"])
        _txb(slide, str(idx + 1).zfill(2),
             x, y, num_w, card_h,
             size=18, bold=True, color=t["white"],
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        # Card body
        _rect(slide, x + num_w, y, card_w - num_w, card_h, t["light"],
              border=t["light"])
        _txb(slide, item,
             x + num_w + Inches(0.15), y, card_w - num_w - Inches(0.15), card_h,
             size=11.5, bold=False, color=t["text"],
             v_anchor=MSO_ANCHOR.MIDDLE)


def _render_executive_summary(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "执行摘要"), "", t)
    _footer(slide, brand_name, num, t)

    headline = data.get("headline", "")
    points   = data.get("points", [])

    # Headline banner
    if headline:
        _rect(slide, MX, CONT_Y + Inches(0.2),
              CONT_W, Inches(0.72), t["primary"])
        _rect(slide, MX, CONT_Y + Inches(0.2),
              Inches(0.09), Inches(0.72), t["accent"])
        _txb(slide, headline,
             MX + Inches(0.24), CONT_Y + Inches(0.2),
             CONT_W - Inches(0.3), Inches(0.72),
             size=14, bold=True, color=t["white"],
             v_anchor=MSO_ANCHOR.MIDDLE)

    # Three insight cards
    n = len(points)
    if not n:
        return
    card_w = CONT_W / max(n, 1) - Inches(0.18)
    card_h = Inches(3.6)
    card_y = CONT_Y + Inches(1.1)

    for i, pt in enumerate(points[:4]):
        cx = MX + i * (card_w + Inches(0.18))
        # Card shadow effect (slightly offset dark rect)
        _rect(slide, cx + Inches(0.05), card_y + Inches(0.06),
              card_w, card_h,
              RGBColor(0xD0, 0xD8, 0xE4))
        # Card
        _rect(slide, cx, card_y, card_w, card_h, t["white"],
              border=t["light"], bw=1.0)
        # Accent top strip
        _rect(slide, cx, card_y, card_w, Inches(0.06), t["accent"])
        # Number
        _txb(slide, f"0{i+1}",
             cx, card_y + Inches(0.15), card_w, Inches(0.75),
             size=36, bold=True, color=t["secondary"],
             align=PP_ALIGN.CENTER)
        # Point text
        _txb(slide, pt,
             cx + Inches(0.15), card_y + Inches(0.95),
             card_w - Inches(0.3), card_h - Inches(1.1),
             size=10.5, color=t["text"], wrap=True)


def _render_section_divider(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    split = Inches(4.8)

    # Left panel (dark)
    _rect(slide, 0, 0, split, SH, t["primary"])
    # Accent stripe
    _rect(slide, split - Inches(0.08), 0, Inches(0.08), SH, t["accent"])
    # Right panel (light)
    _rect(slide, split, 0, SW - split, SH, t["light"])

    # Large section number
    _txb(slide, data.get("number", ""),
         Inches(0.25), SH / 2 - Inches(1.5), split - Inches(0.4), Inches(3),
         size=96, bold=True, color=t["accent"],
         align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # Section title (right side)
    _txb(slide, data.get("title", ""),
         split + Inches(0.5), Inches(2.4), SW - split - Inches(0.8), Inches(1.2),
         size=30, bold=True, color=t["primary"])

    # Horizontal accent line
    _rect(slide, split + Inches(0.5), Inches(3.75),
          SW - split - Inches(0.8), Inches(0.05), t["accent"])

    # Subtitle
    if data.get("subtitle"):
        _txb(slide, data["subtitle"],
             split + Inches(0.5), Inches(3.95),
             SW - split - Inches(0.8), Inches(1.0),
             size=13, color=t["gray"])

    _footer(slide, brand_name, num, t)


def _render_content(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _decor_bg(slide, t)
    _header(slide, data.get("title", ""), data.get("subtitle", ""), t)
    _footer(slide, brand_name, num, t)

    bullets = data.get("bullets", [])
    _bullets_txb(slide, bullets,
                 MX, CONT_Y + Inches(0.3), CONT_W * 0.78, CONT_H - Inches(0.4),
                 size=12, color=t["text"],
                 bullet_color=t["accent"], line_spacing_pt=10)


def _render_two_column(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", ""), "", t)
    _footer(slide, brand_name, num, t)

    col_w = CONT_W / 2 - Inches(0.15)
    left  = data.get("left",  {})
    right = data.get("right", {})
    body_y = CONT_Y + Inches(0.2)

    for i, col in enumerate([left, right]):
        x = MX + i * (col_w + Inches(0.3))
        # Column header
        _rect(slide, x, body_y, col_w, Inches(0.45), t["primary"])
        _txb(slide, col.get("title", ""),
             x + Inches(0.15), body_y, col_w - Inches(0.2), Inches(0.45),
             size=12, bold=True, color=t["white"],
             v_anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        _bullets_txb(slide, col.get("bullets", []),
                     x, body_y + Inches(0.55),
                     col_w, CONT_H - Inches(0.7),
                     size=11, color=t["text"],
                     bullet_color=t["accent"], line_spacing_pt=9)

    # Centre divider
    cx = MX + col_w + Inches(0.15)
    _line(slide, cx, body_y + Inches(0.05), cx, FTR_Y - Inches(0.1),
          t["light"], 1.5)


def _render_swot(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "SWOT战略分析"), "", t)
    _footer(slide, brand_name, num, t)

    quad_defs = [
        ("S", "Strengths  优势",  t["s_clr"],  0, 0),
        ("W", "Weaknesses 劣势",  t["w_clr"],  1, 0),
        ("O", "Opportunities 机会", t["o_clr"], 0, 1),
        ("T", "Threats    威胁",  t["t_clr"],  1, 1),
    ]
    quad_w = CONT_W / 2 - Inches(0.07)
    quad_h = CONT_H / 2 - Inches(0.07)
    hdr_h  = Inches(0.45)

    for key, label, color, col, row in quad_defs:
        x = MX + col * (quad_w + Inches(0.14))
        y = CONT_Y + Inches(0.12) + row * (quad_h + Inches(0.14))

        # Background
        _rect(slide, x, y, quad_w, quad_h, t["light"])
        # Coloured header
        _rect(slide, x, y, quad_w, hdr_h, color)
        _txb(slide, label,
             x + Inches(0.12), y, quad_w - Inches(0.15), hdr_h,
             size=10.5, bold=True, color=t["white"],
             v_anchor=MSO_ANCHOR.MIDDLE)
        # Bullets
        items = data.get(key, [])
        _bullets_txb(slide, items,
                     x + Inches(0.08), y + hdr_h + Inches(0.1),
                     quad_w - Inches(0.16), quad_h - hdr_h - Inches(0.15),
                     size=10, color=t["text"],
                     bullet_color=color, line_spacing_pt=7)


def _render_bar_chart(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", ""), data.get("subtitle", ""), t)
    _footer(slide, brand_name, num, t)

    categories = data.get("categories", [])
    values     = data.get("values", [])
    unit       = data.get("unit", "")
    label      = data.get("label", "数值")

    if not categories or not values:
        return

    cd = ChartData()
    cd.categories = categories
    cd.add_series(label, [float(v) for v in values])

    cx = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        MX, CONT_Y + Inches(0.2),
        CONT_W, CONT_H - Inches(0.35),
        cd,
    )
    chart = cx.chart
    chart.has_legend = False
    chart.has_title  = False

    # Transparent backgrounds
    try:
        chart.chart_area.format.fill.background()
        chart.plot_area.format.fill.background()
    except Exception:
        pass

    # Series colour
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = t["secondary"]

    # Data labels
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size  = Pt(9)
    dl.font.bold  = True
    dl.font.name  = FONT
    dl.font.color.rgb = t["text"]
    if unit:
        dl.number_format = f'0"{unit}"'

    # Axes
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.tick_labels.font.size = Pt(9)
        va.tick_labels.font.name = FONT
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(9)
        ca.tick_labels.font.name = FONT
    except Exception:
        pass


def _render_timeline(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "战略路线图"), "", t)
    _footer(slide, brand_name, num, t)

    phases = data.get("phases", [])
    n = len(phases)
    if not n:
        return

    tl_y   = CONT_Y + CONT_H * 0.44   # timeline axis y
    x0     = MX
    x1     = SW - MX
    col_w  = (x1 - x0) / n

    # Main axis line
    _rect(slide, x0, tl_y - Inches(0.02), x1 - x0, Inches(0.04), t["secondary"])

    # Arrow head
    _rect(slide, x1 - Inches(0.02), tl_y - Inches(0.12),
          Inches(0.04), Inches(0.24), t["secondary"])

    for i, phase in enumerate(phases):
        cx    = x0 + col_w * i + col_w / 2
        color = t["primary"] if i % 2 == 0 else t["secondary"]

        # Timeline node circle
        r = Inches(0.17)
        _oval(slide, cx - r, tl_y - r, r * 2, r * 2, color)

        box_w = col_w - Inches(0.18)
        box_h = Inches(2.1)
        bx    = cx - box_w / 2

        if i % 2 == 0:
            # Above the timeline
            by = tl_y - Inches(0.25) - box_h
            _line(slide, cx, by + box_h, cx, tl_y - r, t["secondary"], 1.2)
        else:
            # Below
            by = tl_y + Inches(0.25)
            _line(slide, cx, tl_y + r, cx, by, t["secondary"], 1.2)

        # Box
        _rect(slide, bx, by, box_w, box_h, t["light"],
              border=color, bw=1.2)
        # Period header
        _rect(slide, bx, by, box_w, Inches(0.38), color)
        _txb(slide, phase.get("period", ""),
             bx + Inches(0.05), by, box_w - Inches(0.1), Inches(0.38),
             size=8.5, bold=True, color=t["white"],
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        # Phase title
        _txb(slide, phase.get("title", ""),
             bx + Inches(0.1), by + Inches(0.42),
             box_w - Inches(0.2), Inches(0.4),
             size=10, bold=True, color=t["primary"])
        # Action items
        items = phase.get("items", [])
        text  = "\n".join(f"• {it}" for it in items[:3])
        _txb(slide, text,
             bx + Inches(0.1), by + Inches(0.86),
             box_w - Inches(0.2), box_h - Inches(1.0),
             size=9, color=t["text"])


def _render_process(prs, data: dict, t: dict, brand_name: str, num: int):
    """
    Professional horizontal process flow.
    Layout: hexagonal step-number badge → tall card with title & description.
    Arrow connectors between cards.
    """
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _decor_bg(slide, t)
    _header(slide, data.get("title", "核心流程"), data.get("subtitle", ""), t)
    _footer(slide, brand_name, num, t)

    steps = data.get("steps", [])
    n = len(steps)
    if not n:
        return

    STEP_COLORS = [
        t["primary"],
        t["secondary"],
        RGBColor(0x1A, 0x78, 0xC2),
        RGBColor(0x0A, 0x93, 0x96),
        RGBColor(0x05, 0xB4, 0x9E),
        RGBColor(0x27, 0xAE, 0x60),
    ]

    # Layout constants
    arrow_w   = Inches(0.42)
    card_h    = CONT_H - Inches(0.15)
    card_w    = (CONT_W - arrow_w * (n - 1)) / n
    card_y    = CONT_Y + Inches(0.08)
    badge_r   = Inches(0.38)   # hexagon badge radius

    for i, step in enumerate(steps):
        x   = MX + i * (card_w + arrow_w)
        clr = STEP_COLORS[i % len(STEP_COLORS)]
        light_clr = _lighten(clr, 0.88)

        # ── Card body (light background) ──
        _rect(slide, x, card_y, card_w, card_h, light_clr, rounded=True)

        # ── Coloured top bar ──
        _rect(slide, x, card_y, card_w, Inches(0.08), clr, rounded=False)

        # ── Step number hexagon badge ──
        hx = x + (card_w - badge_r * 2) / 2
        hy = card_y + Inches(0.22)
        try:
            _hexagon(slide, hx, hy, badge_r * 2, badge_r * 2, clr)
        except Exception:
            _oval(slide, hx, hy, badge_r * 2, badge_r * 2, clr)
        _txb(slide, str(i + 1),
             hx, hy, badge_r * 2, badge_r * 2,
             size=18, bold=True, color=t["white"],
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        # ── Step title ──
        _txb(slide, step.get("title", ""),
             x + Inches(0.12), card_y + Inches(1.08),
             card_w - Inches(0.24), Inches(0.55),
             size=11.5, bold=True, color=clr,
             align=PP_ALIGN.CENTER)

        # ── Divider line under title ──
        _rect(slide, x + card_w * 0.2, card_y + Inches(1.68),
              card_w * 0.6, Inches(0.025), clr)

        # ── Description text ──
        desc = step.get("desc", "")
        if desc:
            _txb(slide, desc,
                 x + Inches(0.14), card_y + Inches(1.82),
                 card_w - Inches(0.28), card_h - Inches(2.0),
                 size=9.5, color=t["text"],
                 align=PP_ALIGN.CENTER, wrap=True)

        # ── Sub-items (bullets inside card) ──
        items = step.get("items", [])
        if items and not desc:
            text = "\n".join(f"• {it}" for it in items[:4])
            _txb(slide, text,
                 x + Inches(0.14), card_y + Inches(1.82),
                 card_w - Inches(0.28), card_h - Inches(2.0),
                 size=9, color=t["text"], wrap=True)

        # ── Arrow connector between cards ──
        if i < n - 1:
            ax   = x + card_w
            ay_c = card_y + card_h / 2
            # Line
            _line(slide, ax, ay_c, ax + arrow_w, ay_c, t["gray"], 1.5)
            # Arrowhead triangle (tiny right-pointing triangle)
            tri_x = ax + arrow_w - Inches(0.14)
            tri_y = ay_c - Inches(0.12)
            _txb(slide, "▶",
                 tri_x, tri_y, Inches(0.18), Inches(0.24),
                 size=11, color=clr,
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)


def _render_flowchart(prs, data: dict, t: dict, brand_name: str, num: int):
    """
    Vertical flowchart with proper shapes:
      oval (start) → rounded-rect (process) → diamond (decision) → oval (end)
    Supports up to ~7 nodes on one slide; beyond that splits into two columns.

    Node format in data["nodes"]:
        {"shape": "oval|rect|diamond", "text": "...", "desc": "...",
         "yes": "Yes label", "no": "No label"}
    """
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _decor_bg(slide, t)
    _header(slide, data.get("title", "业务流程"), data.get("subtitle", ""), t)
    _footer(slide, brand_name, num, t)

    nodes = data.get("nodes", [])
    # Fallback: build nodes from "steps" if "nodes" not present
    if not nodes:
        nodes = []
        steps = data.get("steps", [])
        if steps:
            nodes.append({"shape": "oval",  "text": "开始"})
            for s in steps:
                nodes.append({"shape": "rect",
                               "text": s.get("title", ""),
                               "desc": s.get("desc", "")})
            nodes.append({"shape": "oval",  "text": "完成"})

    n = len(nodes)
    if n == 0:
        return

    # ── Layout ─────────────────────────────────────────────────────
    NODE_W  = Inches(3.6)
    RECT_H  = Inches(0.85)
    OVAL_H  = Inches(0.6)
    DIAM_H  = Inches(0.95)
    GAP     = Inches(0.38)     # vertical gap between nodes
    ARROW_H = Inches(0.28)

    # Calculate total height to center vertically
    def _node_h(nd):
        s = nd.get("shape", "rect")
        return OVAL_H if s == "oval" else (DIAM_H if s == "diamond" else RECT_H)

    total_h = sum(_node_h(nd) for nd in nodes) + GAP * (n - 1)
    start_y = CONT_Y + (CONT_H - total_h) / 2
    cx      = SW / 2

    NODE_COLORS = [t["primary"], t["secondary"],
                   RGBColor(0x1A, 0x78, 0xC2), RGBColor(0x0A, 0x93, 0x96),
                   RGBColor(0x05, 0xB4, 0x9E), RGBColor(0x27, 0xAE, 0x60)]

    y = start_y
    for i, nd in enumerate(nodes):
        shape = nd.get("shape", "rect")
        text  = nd.get("text",  "")
        desc  = nd.get("desc",  "")
        nh    = _node_h(nd)
        nx    = cx - NODE_W / 2
        clr   = NODE_COLORS[i % len(NODE_COLORS)]

        if shape == "oval":
            _oval(slide, nx, y, NODE_W, nh, clr)
            _txb(slide, text, nx, y, NODE_W, nh,
                 size=12, bold=True, color=t["white"],
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

        elif shape == "diamond":
            _diamond(slide, nx, y, NODE_W, nh, _lighten(clr, 0.15),
                     border=clr, bw=2.0)
            _txb(slide, text, nx, y, NODE_W, nh,
                 size=10.5, bold=True, color=t["white"],
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
            # Yes / No labels
            yes_lbl = nd.get("yes", "是")
            no_lbl  = nd.get("no",  "否")
            _txb(slide, f"▼ {yes_lbl}",
                 cx - Inches(0.6), y + nh, Inches(1.2), Inches(0.32),
                 size=8.5, color=t["secondary"],
                 align=PP_ALIGN.CENTER)
            _txb(slide, f"{no_lbl} ▶",
                 nx + NODE_W + Inches(0.05), y + nh / 2 - Inches(0.15),
                 Inches(1.2), Inches(0.32),
                 size=8.5, color=t["accent"])

        else:  # rect (process step)
            _rect(slide, nx, y, NODE_W, nh, clr, rounded=True)
            # Step number badge (small oval on left)
            badge_r = Inches(0.22)
            _oval(slide, nx + Inches(0.12), y + (nh - badge_r * 2) / 2,
                  badge_r * 2, badge_r * 2, t["white"])
            _txb(slide, str(i + 1),
                 nx + Inches(0.12), y + (nh - badge_r * 2) / 2,
                 badge_r * 2, badge_r * 2,
                 size=9, bold=True, color=clr,
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
            # Main text
            _txb(slide, text,
                 nx + Inches(0.62), y, NODE_W - Inches(0.74), nh,
                 size=11, bold=True, color=t["white"],
                 v_anchor=MSO_ANCHOR.MIDDLE)
            # Desc on the right side of the node
            if desc:
                _txb(slide, desc,
                     nx + NODE_W + Inches(0.12), y,
                     SW / 2 - NODE_W / 2 - Inches(0.2), nh,
                     size=9, color=t["gray"],
                     v_anchor=MSO_ANCHOR.MIDDLE)

        # ── Arrow ──────────────────────────────────────────────────
        if i < n - 1:
            arr_x = cx
            arr_y1 = y + nh
            arr_y2 = y + nh + GAP
            _line(slide, arr_x, arr_y1, arr_x, arr_y2, t["secondary"], 1.8)
            _txb(slide, "▼",
                 arr_x - Inches(0.15), arr_y2 - Inches(0.22),
                 Inches(0.3), Inches(0.22),
                 size=9, color=t["secondary"],
                 align=PP_ALIGN.CENTER)

        y += nh + GAP


def _render_visual_content(prs, data: dict, t: dict, brand_name: str, num: int):
    """
    Split layout: left 58% = bullet list, right 42% = large icon + stat card.
    Visually anchors the data with a big number or emoji icon.
    """
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _decor_bg(slide, t)
    _header(slide, data.get("title", ""), data.get("subtitle", ""), t)
    _footer(slide, brand_name, num, t)

    left_w  = CONT_W * 0.58
    right_w = CONT_W * 0.38
    right_x = MX + left_w + CONT_W * 0.04
    body_y  = CONT_Y + Inches(0.25)
    body_h  = CONT_H - Inches(0.35)

    # ── Left: bullets ──────────────────────────────────────────────
    bullets = data.get("bullets", [])
    _bullets_txb(slide, bullets,
                 MX, body_y, left_w, body_h,
                 size=11.5, color=t["text"],
                 bullet_color=t["accent"], line_spacing_pt=12)

    # ── Right: visual card ─────────────────────────────────────────
    card_h = body_h * 0.9
    card_y = body_y + body_h * 0.05
    _rect(slide, right_x, card_y, right_w, card_h, t["light"], rounded=True)
    _rect(slide, right_x, card_y, right_w, Inches(0.07), t["accent"])

    # Large icon / emoji
    icon = data.get("icon", "📊")
    _txb(slide, icon,
         right_x, card_y + Inches(0.35), right_w, Inches(1.4),
         size=64, align=PP_ALIGN.CENTER)

    # Big stat number
    stat = data.get("stat", "")
    if stat:
        _txb(slide, stat,
             right_x, card_y + Inches(1.85), right_w, Inches(0.95),
             size=48, bold=True, color=t["primary"],
             align=PP_ALIGN.CENTER)

    stat_label = data.get("stat_label", "")
    if stat_label:
        _txb(slide, stat_label,
             right_x + Inches(0.1), card_y + Inches(2.85),
             right_w - Inches(0.2), Inches(0.45),
             size=10.5, color=t["gray"],
             align=PP_ALIGN.CENTER)

    # Supporting bullets on right card
    right_bullets = data.get("right_bullets", [])
    if right_bullets:
        rb_y = card_y + Inches(3.45) if stat else card_y + Inches(1.6)
        _bullets_txb(slide, right_bullets,
                     right_x + Inches(0.15), rb_y,
                     right_w - Inches(0.3), card_h - (rb_y - card_y) - Inches(0.1),
                     size=9, color=t["text"],
                     bullet_color=t["accent"], line_spacing_pt=6)


def _render_metrics(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "关键绩效指标"), "", t)
    _footer(slide, brand_name, num, t)

    items = data.get("items", [])
    n = len(items)
    if not n:
        return

    cols = min(n, 4)
    rows = (n + cols - 1) // cols
    gap  = Inches(0.18)
    cw   = (CONT_W - gap * (cols - 1)) / cols
    ch   = (CONT_H - Inches(0.35) - gap * (rows - 1)) / rows

    for idx, item in enumerate(items):
        c = idx % cols
        r = idx // cols
        x = MX + c * (cw + gap)
        y = CONT_Y + Inches(0.2) + r * (ch + gap)

        # Card
        _rect(slide, x, y, cw, ch, t["light"])
        # Top accent line
        _rect(slide, x, y, cw, Inches(0.055), t["accent"])

        # Big value + unit
        val  = str(item.get("value", ""))
        unit = str(item.get("unit", ""))
        _txb(slide, val,
             x, y + Inches(0.18), cw * 0.65, ch - Inches(0.25),
             size=40, bold=True, color=t["primary"],
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        _txb(slide, unit,
             x + cw * 0.65, y + Inches(0.42), cw * 0.35, ch - Inches(0.5),
             size=14, bold=True, color=t["accent"],
             v_anchor=MSO_ANCHOR.MIDDLE)

        # Label
        _txb(slide, item.get("label", ""),
             x + Inches(0.12), y + ch - Inches(0.68),
             cw - Inches(0.24), Inches(0.35),
             size=10.5, bold=True, color=t["text"])

        # Description
        _txb(slide, item.get("desc", ""),
             x + Inches(0.12), y + ch - Inches(0.35),
             cw - Inches(0.24), Inches(0.3),
             size=8.5, color=t["gray"])


def _render_matrix_2x2(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["white"])
    _header(slide, data.get("title", "战略矩阵"), "", t)
    _footer(slide, brand_name, num, t)

    body_y  = CONT_Y + Inches(0.12)
    body_h  = CONT_H - Inches(0.2)
    body_x  = MX + Inches(0.55)   # leave room for y-axis label
    body_w  = CONT_W - Inches(0.6)

    qw = (body_w - Inches(0.12)) / 2
    qh = (body_h - Inches(0.6))  / 2   # leave room for x-axis label

    qcolors = {
        "TL": t["o_clr"],
        "TR": t["s_clr"],
        "BL": t["t_clr"],
        "BR": t["w_clr"],
    }
    qlabels = {"TL": "左上", "TR": "右上", "BL": "左下", "BR": "右下"}
    for q in (data.get("quadrants") or []):
        qlabels[q.get("pos", "")] = q.get("title", "")

    positions = {
        "TL": (body_x,            body_y),
        "TR": (body_x + qw + Inches(0.12), body_y),
        "BL": (body_x,            body_y + qh + Inches(0.12)),
        "BR": (body_x + qw + Inches(0.12), body_y + qh + Inches(0.12)),
    }

    for pos, (qx, qy) in positions.items():
        color = qcolors.get(pos, t["secondary"])
        _rect(slide, qx, qy, qw, qh, t["light"])
        # Coloured header
        _rect(slide, qx, qy, qw, Inches(0.42), color)
        _txb(slide, qlabels.get(pos, ""),
             qx + Inches(0.1), qy, qw - Inches(0.2), Inches(0.42),
             size=11, bold=True, color=t["white"],
             v_anchor=MSO_ANCHOR.MIDDLE)

        # Items
        q_data = next(
            (q for q in (data.get("quadrants") or []) if q.get("pos") == pos),
            {}
        )
        items = q_data.get("items", [])
        _bullets_txb(slide, items,
                     qx + Inches(0.1), qy + Inches(0.5),
                     qw - Inches(0.2), qh - Inches(0.6),
                     size=10, color=t["text"],
                     bullet_color=color, line_spacing_pt=6)

    # Axis labels
    x_axis = data.get("x_axis", "X轴")
    y_axis = data.get("y_axis", "Y轴")
    _txb(slide, x_axis,
         body_x, body_y + body_h - Inches(0.45), body_w, Inches(0.35),
         size=9, color=t["gray"], align=PP_ALIGN.CENTER)
    _txb(slide, y_axis,
         MX, body_y, Inches(0.45), body_h - Inches(0.45),
         size=9, color=t["gray"], align=PP_ALIGN.CENTER)

    # Center cross lines
    cx = body_x + qw + Inches(0.06)
    cy = body_y + qh + Inches(0.06)
    _rect(slide, cx - Inches(0.03), body_y, Inches(0.06),
          body_h - Inches(0.45), t["gray"])
    _rect(slide, body_x, cy - Inches(0.03),
          body_w, Inches(0.06), t["gray"])


def _render_closing(prs, data: dict, t: dict, brand_name: str, num: int):
    slide = _blank(prs)
    _rect(slide, 0, 0, SW, SH, t["primary"])
    _rect(slide, 0, 0, SW, Inches(0.12), t["accent"])
    _rect(slide, 0, SH - Inches(0.1), SW, Inches(0.1), t["accent"])
    _rect(slide, 0, 0, Inches(0.35), SH, t["secondary"])

    _txb(slide, data.get("title", "结论与建议"),
         Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.75),
         size=26, bold=True, color=t["white"])

    _rect(slide, Inches(0.8), Inches(1.7), Inches(8), Inches(0.055), t["accent"])

    # Headline
    headline = data.get("headline", "")
    if headline:
        _txb(slide, headline,
             Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.85),
             size=15, italic=True, color=t["muted"])

    # Numbered action points
    points = data.get("points", [])
    for i, pt in enumerate(points[:5]):
        py = Inches(2.85) + i * Inches(0.85)
        # Number circle
        _oval(slide,
              Inches(0.8), py, Inches(0.5), Inches(0.5), t["accent"])
        _txb(slide, str(i + 1),
             Inches(0.8), py, Inches(0.5), Inches(0.5),
             size=14, bold=True, color=t["white"],
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        # Text
        _txb(slide, pt,
             Inches(1.5), py + Inches(0.05), Inches(10.5), Inches(0.48),
             size=13, color=t["white"],
             v_anchor=MSO_ANCHOR.MIDDLE)

    # Brand watermark
    _txb(slide, brand_name,
         MX, SH - Inches(0.65), Inches(5), Inches(0.4),
         size=9, color=t["muted"])


# ── Dispatch ──────────────────────────────────────────────────────────────────

_RENDERERS = {
    "cover":               _render_cover,
    "agenda":              _render_agenda,
    "executive_summary":   _render_executive_summary,
    "section_divider":     _render_section_divider,
    "content":             _render_content,
    "two_column":          _render_two_column,
    "swot":                _render_swot,
    "bar_chart":           _render_bar_chart,
    "timeline":            _render_timeline,
    "process":             _render_process,
    "flowchart":           _render_flowchart,       # NEW: proper flowchart with diamonds
    "visual_content":      _render_visual_content,  # NEW: text + big icon/stat card
    "metrics":             _render_metrics,
    "matrix_2x2":          _render_matrix_2x2,
    "closing":             _render_closing,
}


def _render_structure(prs, structure: dict, agent_type: str, brand_name: str):
    t = _th(agent_type)
    slides = structure.get("slides", [])
    slide_num = 1
    for sd in slides:
        stype = sd.get("type", "content")
        fn = _RENDERERS.get(stype, _render_content)
        try:
            fn(prs, sd, t, brand_name, slide_num)
        except Exception as exc:
            print(f"[PPTX] Slide '{stype}' error: {exc}")
            # Render a plain fallback slide instead of crashing
            try:
                _render_content(prs, {
                    "title": sd.get("title", stype),
                    "bullets": sd.get("bullets", sd.get("points", [])),
                }, t, brand_name, slide_num)
            except Exception:
                pass
        slide_num += 1


# ── Markdown fallback structure ────────────────────────────────────────────────

def _strip_md(text: str) -> str:
    text = re.sub(r"\*{3}(.+?)\*{3}", r"\1", text)
    text = re.sub(r"\*{2}(.+?)\*{2}", r"\1", text)
    text = re.sub(r"\*(.+?)\*",       r"\1", text)
    text = re.sub(r"`(.+?)`",         r"\1", text)
    return text.strip()


def _parse_markdown_structure(
    brand_name: str,
    content: str,
    agent_type: str,
) -> dict:
    """
    Minimal fallback: parse H2/H3 sections from markdown
    and build a generic slide structure.
    """
    lines = content.split("\n")
    sections: List[dict] = []
    current: Optional[dict] = None

    for line in lines:
        line = line.rstrip()
        if line.startswith("## "):
            if current:
                sections.append(current)
            current = {"title": _strip_md(line[3:]), "bullets": [], "subs": []}
        elif line.startswith("### "):
            if current is not None:
                current["subs"].append({"title": _strip_md(line[4:]), "bullets": []})
        elif line.startswith(("- ", "* ", "• ")) and current:
            txt = _strip_md(line[2:].strip())
            if txt:
                if current["subs"]:
                    current["subs"][-1]["bullets"].append(txt)
                else:
                    current["bullets"].append(txt)
        elif line and not line.startswith("#") and current:
            txt = _strip_md(line)
            if txt and len(txt) < 120:
                if current["subs"]:
                    current["subs"][-1]["bullets"].append(txt)
                elif len(current["bullets"]) < 8:
                    current["bullets"].append(txt)

    if current:
        sections.append(current)

    agent_title = {
        "strategy":   "战略规划方案",
        "brand":      "品牌设计方案",
        "operations": "运营实施方案",
    }.get(agent_type, "专项分析报告")

    slides: List[dict] = [
        {"type": "cover",   "title": f"{brand_name} {agent_title}", "subtitle": "AI 战略咨询报告"},
        {"type": "agenda",  "title": "报告结构",
         "items": [s["title"] for s in sections[:8]]},
    ]

    for i, sec in enumerate(sections):
        # Section divider every 2 sections
        if i % 2 == 0:
            slides.append({
                "type": "section_divider",
                "number": str(i // 2 + 1).zfill(2),
                "title": sec["title"],
                "subtitle": "",
            })

        # Build content slide from section bullets
        all_bullets = sec["bullets"][:]
        for sub in sec["subs"]:
            all_bullets.extend(sub["bullets"])

        slides.append({
            "type": "content",
            "title": sec["title"],
            "bullets": all_bullets[:7],
        })

        # Sub-sections as content slides
        for sub in sec["subs"][:2]:
            if sub["bullets"]:
                slides.append({
                    "type": "content",
                    "title": sub["title"],
                    "subtitle": sec["title"],
                    "bullets": sub["bullets"][:6],
                })

    slides.append({
        "type": "closing",
        "title": "结论与下一步",
        "headline": f"以上是 {brand_name} 的 {agent_title}",
        "points": ["详见正文报告", "欢迎进一步探讨"],
    })

    return {
        "title": f"{brand_name} {agent_title}",
        "subtitle": "AI 战略咨询报告",
        "slides": slides,
    }


# ── Main entry point ──────────────────────────────────────────────────────────

async def generate_pptx(
    task_id: int,
    agent_type: str,
    brand_name: str,
    content: str,
    file_path: str,
    db=None,
) -> None:
    """
    Generate a professional PPTX.

    1. Tries AI-powered structure generation via pptx_agent (if db provided).
    2. Falls back to markdown parsing if AI is unavailable or fails.
    3. Renders to file_path.
    """
    structure = None

    if db is not None:
        try:
            from services.pptx_agent import generate_slide_structure
            structure = await generate_slide_structure(
                brand_name, content, agent_type, db
            )
        except Exception as exc:
            print(f"[PPTX] pptx_agent failed, using fallback: {exc}")

    if structure is None:
        print("[PPTX] Using markdown fallback structure")
        structure = _parse_markdown_structure(brand_name, content, agent_type)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    _render_structure(prs, structure, agent_type, brand_name)
    prs.save(file_path)
    print(f"[PPTX] Saved → {file_path}  ({len(structure['slides'])} slides)")
